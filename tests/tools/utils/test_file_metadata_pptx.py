"""Tests for file_metadata registry with PPTX files (PRE-8 #5439).

Covers:
  * .pptx advertises total_pages + page_number in first_class_params
  * .pptx output conforms to the PRE-1 chunked-read schema (#5432)
  * total_pages reflects the real slide count
  * No-content degrades gracefully to total_pages == 0
  * Corrupt bytes degrade gracefully to total_pages == 0
  * Single-slide bounded read via the loader stays bounded
  * extract_images advertised in extra_params
  * .ppt is no longer in loaders_map (LibreOffice not installed; convert to .pptx)
"""
import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from elitea_sdk.runtime.langchain.document_loaders.EliteAPowerPointLoader import (
    EliteAPowerPointLoader,
)
from elitea_sdk.tools.utils.file_metadata import get_file_metadata


def _make_pptx_bytes(num_slides: int = 3) -> bytes:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    for i in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = f"Slide {i + 1} content"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_metadata_reports_total_pages_and_instruction():
    data = _make_pptx_bytes(3)

    meta = get_file_metadata("deck.pptx", file_content=data, file_size=len(data))

    assert meta["__result_status__"] == "file_metadata"
    assert meta["filename"] == "deck.pptx"
    assert meta["extension"] == ".pptx"
    assert meta["filesize"] == len(data)
    assert meta["read_limits"]["max_output_chars"] == 200000

    assert meta["unit"] == "pages"
    assert meta["total_pages"] == 3

    instr = meta["instruction_for_readFile"]
    assert "page_number" in instr["first_class_params"]
    hint = instr["first_class_params"]["page_number"]
    assert "1-indexed" in hint
    assert "1..3" in hint
    assert "extract_images" in instr["extra_params"]
    assert "page_number" in instr["notes"]


def test_pptx_metadata_no_content():
    meta = get_file_metadata("deck.pptx", file_content=None, file_size=4096)

    assert meta["unit"] == "pages"
    assert meta["total_pages"] == 0
    assert meta["filesize"] == 4096
    assert "page_number" in meta["instruction_for_readFile"]["first_class_params"]


def test_pptx_metadata_corrupt_bytes_degrades():
    meta = get_file_metadata("broken.pptx", file_content=b"not a real pptx",
                             file_size=15)

    assert meta["__result_status__"] == "file_metadata"
    assert meta["total_pages"] == 0
    assert meta["unit"] == "pages"


def test_pptx_metadata_single_slide():
    data = _make_pptx_bytes(1)

    meta = get_file_metadata("single.pptx", file_content=data, file_size=len(data))

    assert meta["total_pages"] == 1
    hint = meta["instruction_for_readFile"]["first_class_params"]["page_number"]
    assert "1..1" in hint


def test_pptx_single_slide_read_is_bounded():
    data = _make_pptx_bytes(3)

    loader = EliteAPowerPointLoader(file_content=data, mode="paged", page_number=2)
    content = loader.get_content()

    # paged mode with page_number returns a single-item list
    assert isinstance(content, list)
    assert len(content) == 1
    assert "Slide 2" in content[0]
    assert "Slide 1" not in content[0]
    assert "Slide 3" not in content[0]


def test_pptx_metadata_extra_params_has_extract_images():
    data = _make_pptx_bytes(2)

    meta = get_file_metadata("img.pptx", file_content=data, file_size=len(data))

    extra = meta["instruction_for_readFile"]["extra_params"]
    assert "extract_images" in extra
    assert "vision" in extra["extract_images"].lower() or "LLM" in extra["extract_images"] or "llm" in extra["extract_images"].lower()


def test_ppt_not_in_loaders_map():
    from elitea_sdk.runtime.langchain.document_loaders.constants import document_loaders_map
    assert ".ppt" not in document_loaders_map, (
        ".ppt requires LibreOffice (not installed); users should convert to .pptx"
    )
